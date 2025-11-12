#!/usr/bin/env python3
"""
Real document processing module
Handles text extraction, chunking, and embeddings
"""

import os
import sys
import re
from typing import List, Dict, Any, Optional
from pathlib import Path
import logging
try:
    # Load root .env so the worker also picks up env vars when started manually
    from dotenv import load_dotenv  # type: ignore
    load_dotenv()
except Exception:
    # If python-dotenv is not installed in this env, continue silently
    pass

# Document processing libraries
import pdfplumber
import pypdf
from docx import Document as DocxDocument
from pptx import Presentation
import pytesseract
from PIL import Image

# AI and embeddings
import openai
from openai import OpenAI
import google.generativeai as genai

# Add parent directory to path
sys.path.append(os.path.join(os.path.dirname(__file__), '..', 'api'))

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class DocumentProcessor:
    """Handles real document processing with text extraction and embeddings"""
    
    def __init__(self):
        # Initialize OpenAI client (will use environment variable)
        self.openai_client = None
        if os.getenv("OPENAI_API_KEY"):
            self.openai_client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
            logger.info("OpenAI client initialized for embeddings")
        
        # Initialize Gemini client
        self.gemini_client = None
        if os.getenv("GEMINI_API_KEY"):
            genai.configure(api_key=os.getenv("GEMINI_API_KEY"))
            self.gemini_client = genai
            logger.info("Gemini client initialized for embeddings")
        
        # Determine which provider to use
        if self.gemini_client:
            self.provider = "gemini"
            logger.info("Using Gemini for embeddings")
            embed_pref = os.getenv("GEMINI_MODEL_EMBED", "text-embedding-004")
            logger.info(f"Gemini embedding model preference: {embed_pref}")
        elif self.openai_client:
            self.provider = "openai"
            logger.info("Using OpenAI for embeddings")
        else:
            self.provider = "mock"
            logger.warning("No API keys found, embeddings will be simulated")
    
    def clean_text_preserve_structure(self, text: str) -> str:
        """Clean text while preserving paragraph structure and meaningful formatting"""
        if not text:
            return ""
        
        # Step 1: Normalize line endings
        text = text.replace('\r\n', '\n').replace('\r', '\n')
        
        # Step 2: Handle hyphenated words at line breaks
        text = re.sub(r'-\n([a-z])', r'\1', text)  # Remove hyphenation
        
        # Step 3: Handle bullet points and numbered lists
        text = re.sub(r'\n\s*([•\-\*\d+\.]\s+)', r'\n\n\1', text)
        
        # Step 4: Preserve paragraph breaks (double newlines)
        # First, mark real paragraph breaks
        text = re.sub(r'\n\s*\n', '||PARAGRAPH_BREAK||', text)
        
        # Step 5: Remove unnecessary single line breaks within paragraphs
        # but keep line breaks before capital letters (likely new sentences/sections)
        text = re.sub(r'\n(?=[a-z])', ' ', text)  # Lowercase continuation -> space
        text = re.sub(r'\n(?=[A-Z][a-z])', '\n', text)  # Capital letter -> keep break
        
        # Step 6: Restore paragraph breaks
        text = text.replace('||PARAGRAPH_BREAK||', '\n\n')
        
        # Step 7: Clean up excessive whitespace
        text = re.sub(r' +', ' ', text)  # Multiple spaces -> single space
        text = re.sub(r'\n{3,}', '\n\n', text)  # Too many newlines -> double
        
        # Step 8: Ensure sentences are properly spaced
        text = re.sub(r'([.!?])([A-Z])', r'\1 \2', text)
        
        return text.strip()
    
    def extract_text_from_pdf(self, file_path: str) -> Dict[str, Any]:
        """Extract text from PDF using pdfplumber and pypdf"""
        text_content = []
        metadata = {"pages": 0, "method": "text_extraction"}
        
        try:
            # Try pdfplumber first (better for text extraction)
            with pdfplumber.open(file_path) as pdf:
                metadata["pages"] = len(pdf.pages)
                
                for page_num, page in enumerate(pdf.pages, 1):
                    text = page.extract_text()
                    if text and text.strip():
                        # Preserve paragraph structure better
                        # Clean up text while preserving meaningful line breaks
                        cleaned_text = self.clean_text_preserve_structure(text)
                        text_content.append({
                            "page": page_num,
                            "text": cleaned_text,
                            "tokens": len(cleaned_text.split())
                        })
                    else:
                        # If no text found, try OCR
                        logger.info(f"No text found on page {page_num}, attempting OCR...")
                        try:
                            # Convert page to image for OCR
                            img = page.to_image()
                            ocr_text = pytesseract.image_to_string(img.original)
                            if ocr_text.strip():
                                cleaned_ocr_text = self.clean_text_preserve_structure(ocr_text)
                                text_content.append({
                                    "page": page_num,
                                    "text": cleaned_ocr_text,
                                    "tokens": len(cleaned_ocr_text.split()),
                                    "method": "ocr"
                                })
                                metadata["method"] = "hybrid"
                        except Exception as ocr_error:
                            logger.warning(f"OCR failed for page {page_num}: {ocr_error}")
                            
        except Exception as e:
            logger.error(f"pdfplumber failed, trying pypdf: {e}")
            
            # Fallback to pypdf
            try:
                with open(file_path, 'rb') as file:
                    pdf_reader = pypdf.PdfReader(file)
                    metadata["pages"] = len(pdf_reader.pages)
                    
                    for page_num, page in enumerate(pdf_reader.pages, 1):
                        text = page.extract_text()
                        if text and text.strip():
                            text_content.append({
                                "page": page_num,
                                "text": text.strip(),
                                "tokens": len(text.split())
                            })
            except Exception as fallback_error:
                logger.error(f"All PDF extraction methods failed: {fallback_error}")
                raise
        
        return {
            "text_content": text_content,
            "metadata": metadata,
            "total_pages": metadata["pages"],
            "total_tokens": sum(item["tokens"] for item in text_content)
        }
    
    def extract_text_from_docx(self, file_path: str) -> Dict[str, Any]:
        """Extract text from DOCX files"""
        try:
            doc = DocxDocument(file_path)
            
            text_content = []
            current_text = []
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    current_text.append(paragraph.text.strip())
            
            # Combine all text
            full_text = "\n".join(current_text)
            
            if full_text.strip():
                text_content.append({
                    "page": 1,
                    "text": full_text.strip(),
                    "tokens": len(full_text.split())
                })
            
            return {
                "text_content": text_content,
                "metadata": {"pages": 1, "method": "docx_extraction"},
                "total_pages": 1,
                "total_tokens": len(full_text.split()) if full_text else 0
            }
            
        except Exception as e:
            logger.error(f"DOCX extraction failed: {e}")
            raise
    
    def extract_text_from_pptx(self, file_path: str) -> Dict[str, Any]:
        """Extract text from PPTX files"""
        try:
            prs = Presentation(file_path)
            
            text_content = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide_text:
                    combined_text = "\n".join(slide_text)
                    text_content.append({
                        "page": slide_num,
                        "text": combined_text,
                        "tokens": len(combined_text.split())
                    })
            
            return {
                "text_content": text_content,
                "metadata": {"pages": len(prs.slides), "method": "pptx_extraction"},
                "total_pages": len(prs.slides),
                "total_tokens": sum(item["tokens"] for item in text_content)
            }
            
        except Exception as e:
            logger.error(f"PPTX extraction failed: {e}")
            raise
    
    def extract_text_from_txt(self, file_path: str) -> Dict[str, Any]:
        """Extract text from TXT files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
            
            if not content.strip():
                # Try different encodings
                for encoding in ['latin-1', 'cp1252', 'iso-8859-1']:
                    try:
                        with open(file_path, 'r', encoding=encoding) as file:
                            content = file.read()
                        if content.strip():
                            break
                    except:
                        continue
            
            text_content = []
            if content.strip():
                text_content.append({
                    "page": 1,
                    "text": content.strip(),
                    "tokens": len(content.split())
                })
            
            return {
                "text_content": text_content,
                "metadata": {"pages": 1, "method": "txt_extraction"},
                "total_pages": 1,
                "total_tokens": len(content.split()) if content else 0
            }
            
        except Exception as e:
            logger.error(f"TXT extraction failed: {e}")
            raise
    
    def extract_text(self, file_path: str, file_type: str) -> Dict[str, Any]:
        """Main text extraction method that routes to appropriate handler"""
        logger.info(f"Extracting text from {file_type} file: {file_path}")
        
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        
        if file_type.lower() == 'pdf':
            return self.extract_text_from_pdf(file_path)
        elif file_type.lower() == 'docx':
            return self.extract_text_from_docx(file_path)
        elif file_type.lower() == 'pptx':
            return self.extract_text_from_pptx(file_path)
        elif file_type.lower() == 'txt':
            return self.extract_text_from_txt(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")
    
    def create_chunks(self, text_content: List[Dict], chunk_size: int = 1000, overlap: int = 200) -> List[Dict[str, Any]]:
        """Create overlapping text chunks for embedding"""
        chunks = []
        chunk_id = 0
        
        for page_data in text_content:
            page_num = page_data["page"]
            text = page_data["text"]
            words = text.split()
            
            # Create chunks for this page
            start = 0
            while start < len(words):
                end = min(start + chunk_size, len(words))
                chunk_text = " ".join(words[start:end])
                
                # Clean up the text
                chunk_text = re.sub(r'\s+', ' ', chunk_text).strip()
                
                if chunk_text:  # Only add non-empty chunks
                    chunks.append({
                        "chunk_id": chunk_id,
                        "page_start": page_num,
                        "page_end": page_num,
                        "text": chunk_text,
                        "tokens": len(chunk_text.split()),
                        "start_word": start,
                        "end_word": end
                    })
                    chunk_id += 1
                
                # Move start position with overlap
                start = max(start + chunk_size - overlap, end)
                if start >= len(words):
                    break
        
        logger.info(f"Created {len(chunks)} chunks")
        return chunks
    
    def generate_embeddings(self, chunks: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        """Generate embeddings for text chunks"""
        if self.provider == "gemini":
            return self._generate_gemini_embeddings(chunks)
        elif self.provider == "openai":
            return self._generate_openai_embeddings(chunks)
        else:
            logger.warning("No API client available, generating mock embeddings")
            # Return mock embeddings for development
            for chunk in chunks:
                chunk["embedding"] = [0.1] * 768  # Standard embedding size for mock
                chunk["embedding_model"] = "mock-embedding"
            return chunks
    
    def _generate_gemini_embeddings(self, chunks: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        """Generate embeddings using Gemini API"""
        try:
            logger.info(f"Generating Gemini embeddings for {len(chunks)} chunks...")
            
            # Try multiple embedding model names
            embedding_models = [
                os.getenv("GEMINI_MODEL_EMBED", "text-embedding-004"),
                "text-embedding-004",
                # Legacy alias kept as last fallback only
                "models/text-embedding-004",
            ]
            
            working_model = None
            # Test which model works
            for model_name in embedding_models:
                try:
                    test_result = self.gemini_client.embed_content(
                        model=model_name,
                        content="test"
                    )
                    working_model = model_name
                    logger.info(f"Using embedding model: {model_name}")
                    break
                except Exception as e:
                    logger.warning(f"Model {model_name} failed: {e}")
                    continue
            
            if not working_model:
                raise Exception("No working Gemini embedding model found")
            
            for i, chunk in enumerate(chunks):
                try:
                    result = self.gemini_client.embed_content(
                        model=working_model,
                        content=chunk["text"]
                    )
                    # Extract embedding values depending on SDK shape
                    values = None
                    try:
                        values = result["embedding"]["values"]
                    except Exception:
                        values = result.get("embedding")

                    if not values:
                        raise ValueError("Empty embedding values returned")

                    chunk["embedding"] = values
                    chunk["embedding_model"] = working_model
                    
                    if (i + 1) % 10 == 0:
                        logger.info(f"Generated embeddings for {i + 1}/{len(chunks)} chunks")
                        
                except Exception as e:
                    logger.error(f"Failed to generate Gemini embedding for chunk {i}: {e}")
                    # Use mock embedding as fallback
                    chunk["embedding"] = [0.1] * 768
                    chunk["embedding_model"] = "fallback-mock"
            
            logger.info("Gemini embedding generation completed")
            return chunks
            
        except Exception as e:
            logger.error(f"Gemini embedding generation failed: {e}")
            # Fallback to mock embeddings
            for chunk in chunks:
                chunk["embedding"] = [0.1] * 768
                chunk["embedding_model"] = "error-fallback"
            return chunks
    
    def _generate_openai_embeddings(self, chunks: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        """Generate embeddings using OpenAI API"""
        try:
            logger.info(f"Generating OpenAI embeddings for {len(chunks)} chunks...")
            
            for i, chunk in enumerate(chunks):
                try:
                    response = self.openai_client.embeddings.create(
                        model="text-embedding-ada-002",
                        input=chunk["text"]
                    )
                    
                    chunk["embedding"] = response.data[0].embedding
                    chunk["embedding_model"] = "text-embedding-ada-002"
                    
                    if (i + 1) % 10 == 0:
                        logger.info(f"Generated embeddings for {i + 1}/{len(chunks)} chunks")
                        
                except Exception as e:
                    logger.error(f"Failed to generate OpenAI embedding for chunk {i}: {e}")
                    # Use mock embedding as fallback
                    chunk["embedding"] = [0.1] * 1536
                    chunk["embedding_model"] = "fallback-mock"
            
            logger.info("OpenAI embedding generation completed")
            return chunks
            
        except Exception as e:
            logger.error(f"OpenAI embedding generation failed: {e}")
            # Fallback to mock embeddings
            for chunk in chunks:
                chunk["embedding"] = [0.1] * 1536
                chunk["embedding_model"] = "error-fallback"
            return chunks
    
    def process_document_full(self, file_path: str, file_type: str, document_id: str) -> Dict[str, Any]:
        """Complete document processing pipeline"""
        logger.info(f"Starting full processing for document {document_id}")
        
        try:
            # Step 1: Extract text
            logger.info("Step 1: Extracting text...")
            extraction_result = self.extract_text(file_path, file_type)
            
            if not extraction_result["text_content"]:
                raise ValueError("No text content could be extracted from document")
            
            # Step 2: Create chunks
            logger.info("Step 2: Creating text chunks...")
            chunks = self.create_chunks(extraction_result["text_content"])
            
            if not chunks:
                raise ValueError("No chunks could be created from extracted text")
            
            # Step 3: Generate embeddings
            logger.info("Step 3: Generating embeddings...")
            chunks_with_embeddings = self.generate_embeddings(chunks)
            
            # Prepare final result
            result = {
                "document_id": document_id,
                "extraction_metadata": extraction_result["metadata"],
                "total_pages": extraction_result["total_pages"],
                "total_tokens": extraction_result["total_tokens"],
                "total_chunks": len(chunks_with_embeddings),
                "chunks": chunks_with_embeddings,
                "processing_status": "success"
            }
            
            logger.info(f"Document {document_id} processed successfully:")
            logger.info(f"  - Pages: {result['total_pages']}")
            logger.info(f"  - Tokens: {result['total_tokens']}")
            logger.info(f"  - Chunks: {result['total_chunks']}")
            
            return result
            
        except Exception as e:
            logger.error(f"Document processing failed for {document_id}: {e}")
            return {
                "document_id": document_id,
                "processing_status": "failed",
                "error": str(e)
            }
